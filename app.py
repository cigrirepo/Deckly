import streamlit as st
import openai
import json
import io
import base64
import re
from PIL import Image
import numpy as np
from sklearn.cluster import KMeans
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import textwrap

# ── App Config ───────────────────────────────────────────────────────────
st.set_page_config(page_title="Deckly – Pro‑grade decks, zero hassle.", layout="wide")
openai.api_key = st.secrets.get("OPENAI_API_KEY", "")

# ── 1) Bullets → Outline + Speaker Notes ────────────────────────────────
@st.cache_data
def bullets_to_outline(bullets: str, tone: str = "Neutral") -> list[dict]:
    """
    Call OpenAI to convert raw bullets into a structured outline with speaker notes.
    Returns a list of sections: {title, points, speaker_notes}.
    """
    system_msg = {
        "role": "system",
        "content": (
            "You are a venture capital analyst and presentation writer. "
            "Receive a list of bullet points about a startup and produce ONLY valid JSON with a single key 'sections'. "
            "The value of 'sections' is a list of objects, each with keys: 'title' (string), 'points' (list of strings), "
            "and 'speaker_notes' (a concise paragraph for presenter notes)."
        )
    }
    user_msg = {
        "role": "user",
        "content": f"""Tone: {tone}
Bullets:
{bullets}"""
    }
    try:
        resp = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[system_msg, user_msg],
            temperature=0.25
        )
        content = resp.choices[0].message.content.strip()
        # Strip Markdown fences if any
        clean = re.sub(r"```json|```", "", content)
        match = re.search(r"\{[\s\S]*\}", clean)
        data = json.loads(match.group()) if match else {}
        sections = data.get("sections", [])
        if not sections:
            raise ValueError("No 'sections' in response")
    except Exception as e:
        st.error(f"Outline generation error: {e}")
        sections = [{
            "title": "Overview",
            "points": bullets.splitlines(),
            "speaker_notes": "Overview of key points."
        }]
    return sections

# ── 2) Palette Extraction ───────────────────────────────────────────────
@st.cache_data
def extract_palette(img_bytes: bytes, k: int = 4) -> list[str]:
    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    arr = np.array(img).reshape(-1, 3)
    km = KMeans(n_clusters=k, n_init="auto").fit(arr)
    colors = km.cluster_centers_.astype(int)
    return [f"#{r:02x}{g:02x}{b:02x}" for r, g, b in colors]

# ── 3) Logo Style Description ─────────────────────────────────────────
@st.cache_data
def describe_logo(img_bytes: bytes) -> str:
    try:
        b64 = base64.b64encode(img_bytes).decode()
        multimodal = [
            {"type": "text", "text": "Describe this logo's style and mood in one sentence."},
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}}
        ]
        resp = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": multimodal}],
            temperature=0.5
        )
        return resp.choices[0].message.content.strip()
    except Exception:
        return ""

# ── 4) PPTX Builder with Accent Bar + Fonts + Notes ─────────────────────
@st.cache_data
def build_deck(sections: list[dict], palette: list[str], logo_bytes: bytes | None = None) -> bytes:
    prs = Presentation()
    slide_w, slide_h = prs.slide_width, prs.slide_height
    hex2rgb = lambda h: tuple(int(h.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
    bg = hex2rgb(palette[0]) if palette else (255, 255, 255)
    accent = hex2rgb(palette[1]) if len(palette) > 1 else (0, 0, 0)

    for idx, sec in enumerate(sections):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        # Background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*bg)
        # Title
        title_shape = slide.shapes.title or slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(1))
        title_shape.text = sec.get('title', '')
        p0 = title_shape.text_frame.paragraphs[0]
        p0.font.size = Pt(42)
        p0.font.color.rgb = RGBColor(*accent)
        # Body
        body_tf = None
        try:
            body_tf = slide.placeholders[1].text_frame
            body_tf.clear()
        except Exception:
            for shp in slide.shapes:
                if getattr(shp, 'has_text_frame', False) and shp != title_shape:
                    body_tf = shp.text_frame
                    body_tf.clear()
                    break
        if body_tf:
            for point in sec.get('points', []):
                p = body_tf.add_paragraph()
                p.text = point
                p.level = 0
                p.font.size = Pt(20)
        # Speaker Notes
        slide.notes_slide.notes_text_frame.text = sec.get('speaker_notes', '')
        # Accent Bar at bottom
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, slide_h - Inches(0.25),
            slide_w, Inches(0.25)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*accent)
        bar.line.color.rgb = RGBColor(*accent)
        # Logo on first slide
        if logo_bytes and idx == 0:
            slide.shapes.add_picture(
                io.BytesIO(logo_bytes),
                slide_w - Inches(1.5), Inches(0.2),
                width=Inches(1.2)
            )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

# ── 5) Executive Summary via LLM → PDF ─────────────────────────────────
@st.cache_data
def render_brief(sections: list[dict]) -> bytes:
    system_msg = {"role": "system", "content": "You are an expert writer of concise one-page executive summaries for investor decks."}
    user_msg = {"role": "user", "content": json.dumps({"sections": sections})}
    try:
        resp = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[system_msg, user_msg],
            temperature=0.4
        )
        summary = resp.choices[0].message.content.strip()
    except Exception:
        summary = "Executive Summary\n\n" + "\n".join(f"- {pt}" for sec in sections for pt in sec.get('points', []))
    # Render PDF
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    margin = 72  # 1 inch
    y = letter[1] - margin
    c.setFont("Helvetica-Bold", 16)
    c.drawString(margin, y, "Executive Summary")
    y -= 36
    c.setFont("Helvetica", 12)
    for paragraph in summary.split("\n\n"):
        for line in textwrap.wrap(paragraph, width=90):
            if y < margin:
                c.showPage()
                y = letter[1] - margin
                c.setFont("Helvetica", 12)
            c.drawString(margin, y, line)
            y -= 18
        y -= 12
    c.save()
    return buf.getvalue()

# ── 6) Streamlit UI ──────────────────────────────────────────────────────
st.title("🛠️ Deckly – Pro‑grade decks, zero hassle.")
with st.sidebar:
    raw = st.text_area("Paste bullets (newline-separated)", height=250)
    tone = st.select_slider("Tone", ["Conservative", "Neutral", "Bold"], value="Neutral")
    logo_file = st.file_uploader("Upload logo (PNG/JPG)", type=["png", "jpg", "jpeg"])
    k = st.slider("Palette size", 2, 8, 4)
if st.button("Generate Deck"):
    if not raw.strip():
        st.warning("Please enter some bullet points to proceed.")
    else:
        sections = bullets_to_outline(raw, tone)
        if logo_file:
            lb = logo_file.read()
            palette = extract_palette(lb, k)
            style_desc = describe_logo(lb)
            if style_desc:
                st.caption(f"Logo style: {style_desc}")
        else:
            palette = ["#2b2b2b", "#007acc", "#585858"]
        deck_bytes = build_deck(sections, palette, lb if logo_file else None)
        pdf_bytes = render_brief(sections)
        with st.expander("🔍 Inspect Sections (Debug)"):
            st.json(sections)
        st.success("Your deck and brief are ready!")
        st.download_button("⬇️ Download PPTX", deck_bytes, "deckly_deck.pptx")
        st.download_button("⬇️ Download PDF", pdf_bytes, "deckly_brief.pdf")
